"""Main SlideRenderer: orchestrates all rendering modules to produce PPTX."""
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..schema.content import SlideContent
from ..schema.design import DesignSpec, DesignTokens, LayoutAssignment
from .oxml_effects import (
    gradient_fill, solid_fill_shape, outer_shadow,
    rounded_corners, inner_border, slide_gradient_background,
)
from .text_render import apply_text_block
from .chart_render import render_chart
from .diagram_render import render_diagram_sync


# Slide dimensions: 16:9 widescreen (EMU)
SLIDE_W = 9144000
SLIDE_H = 5143500


class LayoutGrid:
    """Convert percentage-based positions to EMU coordinates."""

    def __init__(self, tokens: DesignTokens):
        self.tokens = tokens
        self.m = tokens.spacing

    def x(self, pct: float) -> int:
        return int(pct * SLIDE_W)

    def y(self, pct: float) -> int:
        return int(pct * SLIDE_H)

    def w(self, pct: float) -> int:
        return int(pct * SLIDE_W)

    def h(self, pct: float) -> int:
        return int(pct * SLIDE_H)

    @property
    def content_x(self) -> int:
        return self.x(self.m.margin_left_pct)

    @property
    def content_y(self) -> int:
        return self.y(0.15)  # below header band

    @property
    def content_w(self) -> int:
        return self.w(1.0 - self.m.margin_left_pct - self.m.margin_right_pct)

    @property
    def content_h(self) -> int:
        return self.h(1.0 - 0.15 - self.m.margin_bottom_pct - 0.06)


class SlideRenderer:
    """Renders a list of SlideContent objects into a PPTX Presentation."""

    def __init__(self, design_spec: DesignSpec):
        self.spec = design_spec
        self.tokens = design_spec.tokens
        self.grid = LayoutGrid(self.tokens)
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self._blank_layout = self.prs.slide_layouts[6]  # Blank layout

    def render_deck(self, slides: list[SlideContent]) -> Presentation:
        for content in slides:
            layout = self._get_layout(content.index)
            self._render_slide(content, layout)
        return self.prs

    def save(self, path: str | Path) -> None:
        self.prs.save(str(path))

    def _get_layout(self, index: int) -> LayoutAssignment | None:
        for la in self.spec.layout_assignments:
            if la.slide_index == index:
                return la
        return None

    def _render_slide(self, content: SlideContent, layout: LayoutAssignment | None) -> None:
        slide = self.prs.slides.add_slide(self._blank_layout)

        # Background
        self._apply_background(slide, content, layout)

        # Header band (accent stripe)
        self._add_header_stripe(slide)

        # Dispatch by slide type
        render_fn = {
            "title": self._render_title_slide,
            "section": self._render_section_slide,
            "closing": self._render_closing_slide,
        }.get(content.slide_type, self._render_content_slide)

        render_fn(slide, content, layout)

        # Footer: slide number
        self._add_footer(slide, content.index)

        # Speaker notes
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _apply_background(self, slide, content: SlideContent,
                           layout: LayoutAssignment | None) -> None:
        tokens = self.tokens
        variant = content.background_variant

        if variant == "dark" or tokens.use_gradient_background:
            if variant == "dark":
                stops = [(0, tokens.colors.primary), (100, "#0A0F1E")]
            else:
                stops = [(0, tokens.colors.background), (100, tokens.colors.surface)]
            slide_gradient_background(slide, stops, tokens.gradient_angle_deg)
        else:
            # Solid white/surface background via a full-slide rectangle
            from pptx.util import Emu as EmuU
            bg_shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                EmuU(0), EmuU(0), EmuU(SLIDE_W), EmuU(SLIDE_H)
            )
            bg_shape.line.fill.background()
            solid_fill_shape(bg_shape, tokens.colors.background)
            bg_shape.line.width = 0

    def _add_header_stripe(self, slide) -> None:
        """Thin accent stripe at top of slide."""
        from pptx.util import Emu as EmuU
        stripe = slide.shapes.add_shape(
            1, EmuU(0), EmuU(0),
            EmuU(SLIDE_W), EmuU(self.grid.h(0.008))
        )
        solid_fill_shape(stripe, self.tokens.colors.accent)
        stripe.line.width = 0

    def _add_footer(self, slide, slide_index: int) -> None:
        from pptx.util import Emu as EmuU
        footer_y = self.grid.y(0.93)
        footer_h = self.grid.h(0.05)

        # Footer text (optional)
        if self.spec.footer_text:
            txBox = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(footer_y),
                EmuU(self.grid.w(0.6)), EmuU(footer_h)
            )
            tf = txBox.text_frame
            tf.text = self.spec.footer_text
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(self.tokens.typography.caption_size)
            r, g, b = bytes.fromhex(self.tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

        # Slide number
        num_txBox = slide.shapes.add_textbox(
            EmuU(self.grid.x(0.92)), EmuU(footer_y),
            EmuU(self.grid.w(0.06)), EmuU(footer_h)
        )
        tf = num_txBox.text_frame
        tf.text = str(slide_index + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.size = Pt(self.tokens.typography.caption_size)
        r, g, b = bytes.fromhex(self.tokens.colors.text_secondary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _add_heading(self, slide, text: str, y_pct: float = 0.08,
                     size_pt: int | None = None,
                     color_hex: str | None = None) -> None:
        from pptx.util import Emu as EmuU
        size = size_pt or self.tokens.typography.heading_size_h2
        color = color_hex or self.tokens.colors.primary

        txBox = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.y(y_pct)),
            EmuU(self.grid.content_w), EmuU(self.grid.h(0.10))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(size)
        r, g, b = bytes.fromhex(color.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _render_title_slide(self, slide, content: SlideContent,
                             layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU
        tokens = self.tokens

        # Large title
        title_box = slide.shapes.add_textbox(
            EmuU(self.grid.x(0.08)), EmuU(self.grid.y(0.25)),
            EmuU(self.grid.w(0.84)), EmuU(self.grid.h(0.30))
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = content.heading
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(tokens.typography.heading_size_h1)
        r, g, b = bytes.fromhex(tokens.colors.primary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

        # Subtitle
        if content.subheading:
            sub_box = slide.shapes.add_textbox(
                EmuU(self.grid.x(0.10)), EmuU(self.grid.y(0.58)),
                EmuU(self.grid.w(0.80)), EmuU(self.grid.h(0.12))
            )
            tf = sub_box.text_frame
            tf.text = content.subheading
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(tokens.typography.heading_size_h3)
            r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

        # Decorative divider line
        divider = slide.shapes.add_shape(
            1,
            EmuU(self.grid.x(0.30)), EmuU(self.grid.y(0.55)),
            EmuU(self.grid.w(0.40)), EmuU(self.grid.h(0.004))
        )
        solid_fill_shape(divider, tokens.colors.accent)
        divider.line.width = 0

    def _render_section_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU
        tokens = self.tokens

        # Section number badge
        badge = slide.shapes.add_shape(
            1,
            EmuU(self.grid.x(0.06)), EmuU(self.grid.y(0.30)),
            EmuU(self.grid.w(0.08)), EmuU(self.grid.h(0.14))
        )
        solid_fill_shape(badge, tokens.colors.accent)
        rounded_corners(badge, radius_pt=6)
        outer_shadow(badge, blur_pt=8, dist_pt=3, color_hex="#000000", alpha_pct=20)

        tf = badge.text_frame
        tf.text = str(content.index + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(28)
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Section title
        self._add_heading(slide, content.heading, y_pct=0.28,
                          size_pt=tokens.typography.heading_size_h1,
                          color_hex=tokens.colors.primary)

        if content.subheading:
            sub_box = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(self.grid.y(0.52)),
                EmuU(self.grid.content_w), EmuU(self.grid.h(0.15))
            )
            tf = sub_box.text_frame
            tf.text = content.subheading
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(tokens.typography.heading_size_h3)
            r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

    def _render_closing_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        self._render_title_slide(slide, content, layout)

    def _render_content_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU

        # Heading
        self._add_heading(slide, content.heading)

        slide_type = content.slide_type
        la = layout

        if slide_type == "two_column" and la:
            self._render_two_column(slide, content, la)
        elif content.chart is not None:
            self._render_with_chart(slide, content)
        elif content.diagram is not None:
            self._render_with_diagram(slide, content)
        elif content.table is not None:
            self._render_with_table(slide, content)
        else:
            self._render_text_only(slide, content)

    def _render_text_only(self, slide, content: SlideContent) -> None:
        from pptx.util import Emu as EmuU
        if not content.body_blocks:
            return

        txBox = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.content_y),
            EmuU(self.grid.content_w), EmuU(self.grid.content_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for block in content.body_blocks:
            apply_text_block(tf, block, self.tokens)

    def _render_with_chart(self, slide, content: SlideContent) -> None:
        from pptx.util import Emu as EmuU

        # Text on left (if any), chart on right
        has_text = bool(content.body_blocks)
        split = 0.45 if has_text else 0.0

        if has_text:
            txBox = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(self.grid.content_y),
                EmuU(self.grid.w(split - 0.03)), EmuU(self.grid.content_h)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for block in content.body_blocks:
                apply_text_block(tf, block, self.tokens)

        chart_x = self.grid.x(split + self.tokens.spacing.margin_left_pct if has_text else 0.06)
        chart_w = self.grid.w(1.0 - split - self.tokens.spacing.margin_right_pct - (
            self.tokens.spacing.margin_left_pct if has_text else 0.06))
        chart_h = self.grid.content_h

        chart_png = render_chart(content.chart, self.tokens,
                                  width_px=int(chart_w / 9144 * 96),
                                  height_px=int(chart_h / 5143500 * 540 * 2))
        slide.shapes.add_picture(chart_png,
                                  EmuU(chart_x), EmuU(self.grid.content_y),
                                  EmuU(chart_w), EmuU(chart_h))

    def _render_with_diagram(self, slide, content: SlideContent) -> None:
        from pptx.util import Emu as EmuU

        diag_png = render_diagram_sync(content.diagram.mermaid,
                                        content.diagram.caption or "diagram")

        diag_x = self.grid.content_x
        diag_y = self.grid.content_y
        diag_w = self.grid.content_w
        diag_h = self.grid.content_h

        if content.body_blocks:
            # Text top, diagram bottom
            text_h = self.grid.h(0.20)
            txBox = slide.shapes.add_textbox(
                EmuU(diag_x), EmuU(diag_y),
                EmuU(diag_w), EmuU(text_h)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for block in content.body_blocks:
                apply_text_block(tf, block, self.tokens)
            diag_y += text_h + self.grid.h(0.02)
            diag_h -= text_h + self.grid.h(0.02)

        slide.shapes.add_picture(diag_png,
                                  EmuU(diag_x), EmuU(diag_y),
                                  EmuU(diag_w), EmuU(diag_h))

    def _render_with_table(self, slide, content: SlideContent) -> None:
        from pptx.util import Emu as EmuU
        from pptx.util import Inches

        spec = content.table
        rows = len(spec.rows) + 1  # +1 for header
        cols = len(spec.headers)
        tokens = self.tokens

        tbl = slide.shapes.add_table(
            rows, cols,
            EmuU(self.grid.content_x), EmuU(self.grid.content_y),
            EmuU(self.grid.content_w), EmuU(self.grid.content_h)
        ).table

        # Header row
        for ci, header in enumerate(spec.headers):
            cell = tbl.cell(0, ci)
            cell.text = header
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.bold = True
            run.font.size = Pt(tokens.typography.body_size)
            r, g, b = bytes.fromhex(tokens.colors.background.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)
            from pptx.oxml.ns import qn
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = tcPr.find(qn("a:solidFill"))
            if solidFill is not None:
                tcPr.remove(solidFill)
            from lxml import etree
            sf = etree.SubElement(tcPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
            etree.SubElement(sf, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
                             val=tokens.colors.primary.lstrip("#"))

        # Data rows
        for ri, row in enumerate(spec.rows):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = cell_text
                para = cell.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()
                run.font.size = Pt(tokens.typography.body_size - 1)
                r, g, b = bytes.fromhex(tokens.colors.text_primary.lstrip("#"))
                run.font.color.rgb = RGBColor(r, g, b)
                if spec.highlight_column is not None and ci == spec.highlight_column:
                    run.font.bold = True

    def _render_two_column(self, slide, content: SlideContent,
                            layout: LayoutAssignment) -> None:
        from pptx.util import Emu as EmuU

        split = layout.split_ratio
        gutter = self.tokens.spacing.gutter_pct
        left_w = self.grid.w(split - gutter / 2)
        right_x = self.grid.x(self.tokens.spacing.margin_left_pct + split + gutter / 2)
        right_w = self.grid.w(1.0 - self.tokens.spacing.margin_left_pct
                               - self.tokens.spacing.margin_right_pct - split - gutter / 2)

        # Left: text blocks
        if content.body_blocks:
            txBox = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(self.grid.content_y),
                EmuU(left_w), EmuU(self.grid.content_h)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for block in content.body_blocks:
                apply_text_block(tf, block, self.tokens)

        # Right: chart, diagram, or table
        if content.chart:
            chart_png = render_chart(content.chart, self.tokens)
            slide.shapes.add_picture(chart_png,
                                      EmuU(right_x), EmuU(self.grid.content_y),
                                      EmuU(right_w), EmuU(self.grid.content_h))
        elif content.diagram:
            diag_png = render_diagram_sync(content.diagram.mermaid,
                                            content.diagram.caption or "diagram")
            slide.shapes.add_picture(diag_png,
                                      EmuU(right_x), EmuU(self.grid.content_y),
                                      EmuU(right_w), EmuU(self.grid.content_h))
        elif content.table:
            # Re-render table in right zone
            from pptx.util import Emu as EmuU
            spec = content.table
            rows = len(spec.rows) + 1
            cols = len(spec.headers)
            slide.shapes.add_table(rows, cols,
                                    EmuU(right_x), EmuU(self.grid.content_y),
                                    EmuU(right_w), EmuU(self.grid.content_h))
