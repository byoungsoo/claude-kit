"""Main SlideRenderer: orchestrates all rendering modules to produce PPTX."""
import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..schema.content import SlideContent
from ..schema.design import DesignSpec, DesignTokens, LayoutAssignment
from .oxml_effects import (
    gradient_fill, solid_fill_shape, outer_shadow,
    rounded_corners, inner_border, slide_gradient_background,
)
from .text_render import apply_text_block
from .chart_render import render_chart
from .diagram_render import render_svg, render_aws_diagram
from .layout_templates import get_template, SLIDE_W as _LT_W, SLIDE_H as _LT_H


# Slide dimensions: 16:9 widescreen (EMU)
SLIDE_W = 9144000
SLIDE_H = 5143500


class LayoutGrid:
    """Convert percentage-based positions to EMU coordinates."""

    def __init__(self, tokens: DesignTokens):
        self.tokens = tokens
        self.m = tokens.spacing

    def x(self, pct: float) -> int:
        return int(pct * SLIDE_W)

    def y(self, pct: float) -> int:
        return int(pct * SLIDE_H)

    def w(self, pct: float) -> int:
        return int(pct * SLIDE_W)

    def h(self, pct: float) -> int:
        return int(pct * SLIDE_H)

    @property
    def content_x(self) -> int:
        return self.x(self.m.margin_left_pct)

    @property
    def content_y(self) -> int:
        # heading(0.08~0.20) + subheading(~0.29) 아래부터 시작
        return self.y(0.31)

    @property
    def content_w(self) -> int:
        return self.w(1.0 - self.m.margin_left_pct - self.m.margin_right_pct)

    @property
    def content_h(self) -> int:
        return self.h(1.0 - 0.27 - self.m.margin_bottom_pct - 0.06)


class SlideRenderer:
    """Renders a list of SlideContent objects into a PPTX Presentation."""

    def __init__(self, design_spec: DesignSpec):
        self.spec = design_spec
        self.tokens = design_spec.tokens
        self.grid = LayoutGrid(self.tokens)
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self._blank_layout = self.prs.slide_layouts[6]  # Blank layout
        self._clear_master_background()

    def _clear_master_background(self) -> None:
        """Remove gradient/theme background from slide master so slides start truly blank."""
        from pptx.oxml.ns import qn
        from lxml import etree
        for master in self.prs.slide_masters:
            bg = master._element.find(qn("p:bg"))
            if bg is not None:
                master._element.remove(bg)
            for layout in master.slide_layouts:
                bg = layout._element.find(qn("p:bg"))
                if bg is not None:
                    layout._element.remove(bg)

    def render_deck(self, slides: list[SlideContent]) -> Presentation:
        for content in slides:
            layout = self._get_layout(content.index)
            self._render_slide(content, layout)
        return self.prs

    def save(self, path: str | Path) -> None:
        self.prs.save(str(path))

    def _get_layout(self, index: int) -> LayoutAssignment | None:
        for la in self.spec.layout_assignments:
            if la.slide_index == index:
                return la
        return None

    def _render_slide(self, content: SlideContent, layout: LayoutAssignment | None) -> None:
        slide = self.prs.slides.add_slide(self._blank_layout)

        # Background
        self._apply_background(slide, content, layout)

        # Header band (accent stripe)
        self._add_header_stripe(slide)

        # Dispatch by slide type
        render_fn = {
            "title": self._render_title_slide,
            "section": self._render_section_slide,
            "closing": self._render_closing_slide,
        }.get(content.slide_type, self._render_content_slide)

        render_fn(slide, content, layout)

        # Footer: slide number
        self._add_footer(slide, content.index)

        # Speaker notes
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _apply_background(self, slide, content: SlideContent,
                           layout: LayoutAssignment | None) -> None:
        """Set slide background color via the slide's <p:bg> XML element."""
        from pptx.oxml.ns import qn as _qn
        from lxml import etree
        tokens = self.tokens
        color_hex = tokens.colors.background.lstrip("#")

        # Build <p:bg><p:bgPr><a:solidFill>...
        spTree = slide.shapes._spTree
        sp_parent = spTree.getparent()  # <p:cSld>
        cSld = slide._element.find(_qn("p:cSld"))

        # Remove existing bg if any
        existing = cSld.find(_qn("p:bg"))
        if existing is not None:
            cSld.remove(existing)

        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bg_el = etree.SubElement(cSld, f"{{{NS_P}}}bg")
        bgPr = etree.SubElement(bg_el, f"{{{NS_P}}}bgPr")
        solidFill = etree.SubElement(bgPr, f"{{{NS_A}}}solidFill")
        etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr", val=color_hex)
        etree.SubElement(bgPr, f"{{{NS_A}}}effectLst")

        # Move bg to be the first child of cSld
        cSld.remove(bg_el)
        cSld.insert(0, bg_el)

    def _add_header_stripe(self, slide) -> None:
        pass

    def _add_footer(self, slide, slide_index: int) -> None:
        from pptx.util import Emu as EmuU
        footer_y = self.grid.y(0.93)
        footer_h = self.grid.h(0.05)

        # Footer text (optional)
        if self.spec.footer_text:
            txBox = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(footer_y),
                EmuU(self.grid.w(0.6)), EmuU(footer_h)
            )
            tf = txBox.text_frame
            tf.text = self.spec.footer_text
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.name = self.tokens.typography.body_font
            run.font.size = Pt(self.tokens.typography.caption_size)
            r, g, b = bytes.fromhex(self.tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

        # Slide number
        num_txBox = slide.shapes.add_textbox(
            EmuU(self.grid.x(0.92)), EmuU(footer_y),
            EmuU(self.grid.w(0.06)), EmuU(footer_h)
        )
        tf = num_txBox.text_frame
        tf.text = str(slide_index + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.name = self.tokens.typography.body_font
        run.font.size = Pt(self.tokens.typography.caption_size)
        r, g, b = bytes.fromhex(self.tokens.colors.text_secondary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _add_heading(self, slide, text: str, y_pct: float = 0.08,
                     size_pt: int | None = None,
                     color_hex: str | None = None) -> None:
        from pptx.util import Emu as EmuU
        size = size_pt or self.tokens.typography.heading_size_h2
        color = color_hex or self.tokens.colors.primary

        txBox = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.y(y_pct)),
            EmuU(self.grid.content_w), EmuU(self.grid.h(0.14))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.name = self.tokens.typography.heading_font
        run.font.bold = True
        run.font.size = Pt(size)
        r, g, b = bytes.fromhex(color.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _add_subheading(self, slide, text: str) -> None:
        from pptx.util import Emu as EmuU
        tokens = self.tokens
        txBox = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.y(0.20)),
            EmuU(self.grid.content_w), EmuU(self.grid.h(0.09))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.name = tokens.typography.body_font
        run.font.size = Pt(tokens.typography.heading_size_h3)
        r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _render_title_slide(self, slide, content: SlideContent,
                             layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU
        tokens = self.tokens

        # Title — 세로 중앙 정렬
        title_box = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.y(0.30)),
            EmuU(self.grid.content_w), EmuU(self.grid.h(0.35))
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = content.heading
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = tokens.typography.heading_font
        run.font.bold = True
        run.font.size = Pt(tokens.typography.heading_size_h1 + 4)
        r, g, b = bytes.fromhex(tokens.colors.primary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

        # Subtitle
        if content.subheading:
            sub_box = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(self.grid.y(0.66)),
                EmuU(self.grid.content_w), EmuU(self.grid.h(0.12))
            )
            tf = sub_box.text_frame
            tf.text = content.subheading
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = tokens.typography.body_font
            run.font.size = Pt(tokens.typography.heading_size_h3 + 2)
            r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

    def _render_section_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU
        tokens = self.tokens

        # 섹션 번호
        num_box = slide.shapes.add_textbox(
            EmuU(self.grid.content_x), EmuU(self.grid.y(0.25)),
            EmuU(self.grid.content_w), EmuU(self.grid.h(0.15))
        )
        tf = num_box.text_frame
        tf.text = f"Section {content.index:02d}"
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.name = tokens.typography.body_font
        run.font.size = Pt(tokens.typography.body_size)
        r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
        run.font.color.rgb = RGBColor(r, g, b)

        # 섹션 제목
        self._add_heading(slide, content.heading, y_pct=0.38,
                          size_pt=tokens.typography.heading_size_h1 + 2,
                          color_hex=tokens.colors.primary)

        if content.subheading:
            sub_box = slide.shapes.add_textbox(
                EmuU(self.grid.content_x), EmuU(self.grid.y(0.60)),
                EmuU(self.grid.content_w), EmuU(self.grid.h(0.15))
            )
            tf = sub_box.text_frame
            tf.text = content.subheading
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.name = tokens.typography.body_font
            run.font.size = Pt(tokens.typography.heading_size_h3)
            r, g, b = bytes.fromhex(tokens.colors.text_secondary.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)

    def _render_closing_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        self._render_title_slide(slide, content, layout)

    # ------------------------------------------------------------------ #
    #  Helpers: aspect-ratio-safe picture placement
    # ------------------------------------------------------------------ #

    @staticmethod
    def _svg_aspect(svg_str: str) -> float | None:
        """Return width/height from SVG viewBox or width/height attrs, or None."""
        m = re.search(r'viewBox=["\'][\d.]+\s+[\d.]+\s+([\d.]+)\s+([\d.]+)["\']', svg_str)
        if m:
            w, h = float(m.group(1)), float(m.group(2))
            return w / h if h else None
        mw = re.search(r'\bwidth=["\'](\d+\.?\d*)', svg_str)
        mh = re.search(r'\bheight=["\'](\d+\.?\d*)', svg_str)
        if mw and mh:
            w, h = float(mw.group(1)), float(mh.group(1))
            return w / h if h else None
        return None

    def _place_picture(self, slide, png_buf: io.BytesIO,
                       zone: tuple[int, int, int, int],
                       aspect: float | None = None) -> None:
        """Place a picture inside a zone while preserving aspect ratio (letterbox)."""
        from pptx.util import Emu as EmuU
        zx, zy, zw, zh = zone
        if aspect is not None and aspect > 0:
            zone_aspect = zw / zh
            if aspect > zone_aspect:
                # wider than zone → fit width, letterbox top/bottom
                pw = zw
                ph = int(zw / aspect)
                px = zx
                py = zy + (zh - ph) // 2
            else:
                # taller than zone → fit height, letterbox left/right
                ph = zh
                pw = int(zh * aspect)
                py = zy
                px = zx + (zw - pw) // 2
        else:
            px, py, pw, ph = zx, zy, zw, zh
        slide.shapes.add_picture(png_buf, EmuU(px), EmuU(py), EmuU(pw), EmuU(ph))

    # ------------------------------------------------------------------ #
    #  Template-based content rendering
    # ------------------------------------------------------------------ #

    def _render_content_slide(self, slide, content: SlideContent,
                               layout: LayoutAssignment | None) -> None:
        from pptx.util import Emu as EmuU

        self._add_heading(slide, content.heading)
        if content.subheading:
            self._add_subheading(slide, content.subheading)

        # Determine template: prefer layout_template, fall back to legacy logic
        tpl_name = content.layout_template or self._infer_template(content, layout)
        tpl = get_template(tpl_name)

        if tpl_name == "two_column_text":
            self._render_two_column_text(slide, content, tpl)
        elif tpl_name == "text_only":
            self._render_text_only(slide, content)
        else:
            self._render_by_template(slide, content, tpl, tpl_name)

    def _infer_template(self, content: SlideContent,
                        layout: LayoutAssignment | None) -> str:
        """Legacy fallback: infer template from slide_type and content fields."""
        if content.slide_type == "two_column" and layout:
            return "two_column_text"
        if content.chart or content.diagram or content.table:
            if content.body_blocks:
                return "text_left_viz_right"
            return "viz_only"
        return "text_only"

    def _render_by_template(self, slide, content: SlideContent,
                             tpl: dict, tpl_name: str) -> None:
        from pptx.util import Emu as EmuU

        # Render body text if template has a "text" zone
        if "text" in tpl and content.body_blocks:
            zx, zy, zw, zh = tpl["text"]
            txBox = slide.shapes.add_textbox(
                EmuU(zx), EmuU(zy), EmuU(zw), EmuU(zh)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for block in content.body_blocks:
                apply_text_block(tf, block, self.tokens)

        # Render visualization if template has a "viz" zone
        if "viz" not in tpl:
            return

        viz_zone = tpl["viz"]
        zx, zy, zw, zh = viz_zone

        if content.chart:
            self._add_card(slide, zx, zy, zw, zh)
            chart_png = render_chart(
                content.chart, self.tokens,
                width_px=max(100, int(zw / 914400 * 96)),
                height_px=max(100, int(zh / 914400 * 96)),
            )
            # Charts are rendered to exact pixel size, no aspect correction needed
            slide.shapes.add_picture(chart_png, EmuU(zx), EmuU(zy), EmuU(zw), EmuU(zh))

        elif content.diagram:
            if content.diagram.aws_diagram:
                diag_png = render_aws_diagram(content.diagram.aws_diagram)
                aspect = None  # unknown aspect for binary PNG
            else:
                svg_str = content.diagram.svg or ""
                aspect = self._svg_aspect(svg_str)
                diag_png = render_svg(svg_str)

            self._add_card(slide, zx, zy, zw, zh)
            self._place_picture(slide, diag_png, viz_zone, aspect)

        elif content.table:
            self._add_card(slide, zx, zy, zw, zh)
            self._render_table_in_zone(slide, content.table, zx, zy, zw, zh)

    def _render_text_only(self, slide, content: SlideContent) -> None:
        from pptx.util import Emu as EmuU
        if not content.body_blocks:
            return
        tpl = get_template("text_only")
        zx, zy, zw, zh = tpl["text"]
        txBox = slide.shapes.add_textbox(EmuU(zx), EmuU(zy), EmuU(zw), EmuU(zh))
        tf = txBox.text_frame
        tf.word_wrap = True
        for block in content.body_blocks:
            apply_text_block(tf, block, self.tokens)

    def _add_card(self, slide, x: int, y: int, w: int, h: int) -> None:
        """Render a subtle card background behind content."""
        from pptx.util import Emu as EmuU
        pad = self.grid.w(0.008)
        card = slide.shapes.add_shape(
            1,
            EmuU(x - pad), EmuU(y - pad),
            EmuU(w + pad * 2), EmuU(h + pad * 2)
        )
        solid_fill_shape(card, self.tokens.colors.surface)
        card.line.fill.solid()
        r, g, b = bytes.fromhex(self.tokens.colors.border.lstrip("#"))
        card.line.color.rgb = RGBColor(r, g, b)
        from pptx.util import Pt as PtU
        card.line.width = PtU(0.5)

    def _render_table_in_zone(self, slide, spec, zx: int, zy: int, zw: int, zh: int) -> None:
        from pptx.util import Emu as EmuU
        from lxml import etree
        from pptx.oxml.ns import qn
        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tokens = self.tokens
        rows = len(spec.rows) + 1
        cols = len(spec.headers)

        # 행이 많을수록 폰트 자동 축소 (zone 높이 기준)
        # 각 행 최소 높이: 914400 EMU = 1인치 ≈ 72pt → 행당 약 20pt 이상 필요
        row_h_emu = zh / rows
        base_body = tokens.typography.body_size
        # row 높이(pt 환산)의 55% 이하로 폰트 제한
        max_font = max(8, int(row_h_emu / 914400 * 72 * 0.55))
        body_size = min(base_body - 1, max_font)
        header_size = min(base_body, max_font + 1)

        tbl = slide.shapes.add_table(
            rows, cols, EmuU(zx), EmuU(zy), EmuU(zw), EmuU(zh)
        ).table
        for ci, header in enumerate(spec.headers):
            cell = tbl.cell(0, ci)
            cell.text = header
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = tokens.typography.heading_font
            run.font.bold = True
            run.font.size = Pt(header_size)
            r, g, b = bytes.fromhex(tokens.colors.background.lstrip("#"))
            run.font.color.rgb = RGBColor(r, g, b)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            sf_existing = tcPr.find(qn("a:solidFill"))
            if sf_existing is not None:
                tcPr.remove(sf_existing)
            sf = etree.SubElement(tcPr, f"{{{NS}}}solidFill")
            etree.SubElement(sf, f"{{{NS}}}srgbClr", val=tokens.colors.primary.lstrip("#"))
        for ri, row_data in enumerate(spec.rows):
            row_color = tokens.colors.surface if ri % 2 == 0 else tokens.colors.background
            for ci, cell_text in enumerate(row_data):
                cell = tbl.cell(ri + 1, ci)
                cell.text = cell_text
                para = cell.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()
                run.font.name = tokens.typography.body_font
                run.font.size = Pt(body_size)
                r, g, b = bytes.fromhex(tokens.colors.text_primary.lstrip("#"))
                run.font.color.rgb = RGBColor(r, g, b)
                if spec.highlight_column is not None and ci == spec.highlight_column:
                    run.font.bold = True
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                sf_existing = tcPr.find(qn("a:solidFill"))
                if sf_existing is not None:
                    tcPr.remove(sf_existing)
                sf = etree.SubElement(tcPr, f"{{{NS}}}solidFill")
                etree.SubElement(sf, f"{{{NS}}}srgbClr", val=row_color.lstrip("#"))

    def _render_two_column_text(self, slide, content: SlideContent,
                                 tpl: dict) -> None:
        """Render body_blocks split across two text zones."""
        from pptx.util import Emu as EmuU
        blocks = content.body_blocks
        mid = (len(blocks) + 1) // 2
        left_blocks = blocks[:mid]
        right_blocks = blocks[mid:]

        for zone_key, zone_blocks in [("left", left_blocks), ("right", right_blocks)]:
            if zone_key not in tpl or not zone_blocks:
                continue
            zx, zy, zw, zh = tpl[zone_key]
            txBox = slide.shapes.add_textbox(EmuU(zx), EmuU(zy), EmuU(zw), EmuU(zh))
            tf = txBox.text_frame
            tf.word_wrap = True
            for block in zone_blocks:
                apply_text_block(tf, block, self.tokens)
